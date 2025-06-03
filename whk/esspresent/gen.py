from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# Create presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Exxon Valdez Oil Spill: Anatomy of an Environmental Disaster"
subtitle.text = "Corporate Negligence, Ecological Devastation & Lessons Unlearned\n\nPresentation Date"
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.size = Pt(36)
subtitle.text_frame.paragraphs[0].font.size = Pt(20)

# Slide 2: Introduction/Background
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Introduction/Background"
content = slide.shapes.placeholders[1]
tf = content.text_frame
tf.text = "Key Facts:"
p = tf.add_paragraph()
p.text = "• March 24, 1989: Exxon Valdez supertanker strikes Bligh Reef"
p = tf.add_paragraph()
p.text = "• 11 million gallons spilled (≈125 Olympic swimming pools)"
p = tf.add_paragraph()
p.text = "• Contaminated 1,300+ miles of Alaskan coastline"
p = tf.add_paragraph()
p.text = "• Largest U.S. spill until Deepwater Horizon (2010)"
p = tf.add_paragraph()
p.text = "\nImmediate Consequences:"
p = tf.add_paragraph()
p.text = "• 4,000 km² oil slick within 72 hours"
p = tf.add_paragraph()
p.text = "• Decades-long ecological crisis"

# Add speaker notes
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = ("SPEAKER NOTES:\n"
                                     "- Emphasize scale: Equivalent to emptying 125 Olympic swimming pools of crude oil\n"
                                     "- Show timeline: 1989-1992 key events\n"
                                     "VISUAL: Aerial photo of oil slick")

# Slide 3: Location
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Location: Prince William Sound, Alaska"
content = slide.shapes.placeholders[1]
tf = content.text_frame
tf.text = "Ecological Significance:"
p = tf.add_paragraph()
p.text = "• Critical habitat for salmon, orcas, seabirds, sea otters"
p = tf.add_paragraph()
p.text = "• Home to Chugach Indigenous communities"
p = tf.add_paragraph()
p.text = "\nGeographic Vulnerability:"
p = tf.add_paragraph()
p.text = "• Cold water slows natural oil degradation"
p = tf.add_paragraph()
p.text = "• Fjords trap oil in sediments"
p = tf.add_paragraph()
p.text = "\nLong-term Contamination:"
p = tf.add_paragraph()
p.text = "• 25,000+ gallons remained after 20 years"
p = tf.add_paragraph()
p.text = "• Toxins still detectable today"

# Speaker notes
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = ("SPEAKER NOTES:\n"
                                     "- Highlight biodiversity: 200+ species impacted\n"
                                     "- Explain why cold environments are particularly vulnerable to oil persistence\n"
                                     "VISUAL: Map highlighting spill zone + ecosystem hotspots")

# Slide 4: Source/Cause of Problem
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Source/Cause of the Problem"
content = slide.shapes.placeholders[1]
tf = content.text_frame
tf.text = "Direct Causes:"
p = tf.add_paragraph()
p.text = "• Intoxicated captain asleep below deck"
p = tf.add_paragraph()
p.text = "• Untrained third mate navigating"
p = tf.add_paragraph()
p.text = "• Radar broken for 1+ year (cost-cutting measure)"
p = tf.add_paragraph()
p.text = "\nSystemic Failures:"
p = tf.add_paragraph()
p.text = "• Fraudulent spill-response plans:"
p = tf.add_paragraph()
p.text = "  - 'Ghost teams' (non-existent responders)"
p = tf.add_paragraph()
p.text = "  - Falsified equipment logs"
p = tf.add_paragraph()
p.text = "• Regulatory failures:"
p = tf.add_paragraph()
p.text = "  - Ignored 1971 safety pledges to Indigenous communities"

# Speaker notes
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = ("SPEAKER NOTES:\n"
                                     "- Emphasize corporate negligence over individual blame\n"
                                     "- Mention Exxon saved $22,000/year by not repairing radar\n"
                                     "VISUAL: Infographic showing human/system failure chain")

# Slide 5: Impact on Environment
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Impact on Environment"
content = slide.shapes.placeholders[1]
tf = content.text_frame
tf.text = "Immediate Devastation:"
p = tf.add_paragraph()
p.text = "• 250,000 seabirds killed"
p = tf.add_paragraph()
p.text = "• 2,800 sea otters perished"
p = tf.add_paragraph()
p.text = "• 300 seals dead"
p = tf.add_paragraph()
p.text = "\nLong-term Consequences:"
p = tf.add_paragraph()
p.text = "• Oil persists in sediments (>30 years)"
p = tf.add_paragraph()
p.text = "• Trophic cascade:"
p = tf.add_paragraph()
p.text = "  Otter decline → Urchin explosion → Kelp forest collapse"
p = tf.add_paragraph()
p.text = "• Indigenous communities:"
p = tf.add_paragraph()
p.text = "  1/3 fishermen died before receiving compensation"

# Speaker notes
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = ("SPEAKER NOTES:\n"
                                     "- Quantify impacts: 28% of otters, 9% of birds in area killed\n"
                                     "- Mention PAH concentrations caused 80% amphipod mortality\n"
                                     "VISUAL: Side-by-side wildlife photos (oiled vs healthy)")

# Slide 6: Management Strategy
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Management Strategy"
content = slide.shapes.placeholders[1]
tf = content.text_frame

# Create table
rows, cols = 3, 3
left = Inches(0.5)
top = Inches(2)
width = Inches(9)
height = Inches(4)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Table headers
table.cell(0, 0).text = "Failed Methods"
table.cell(0, 1).text = "Effective Solutions"
table.cell(0, 2).text = "Policy Response"

# Table content
table.cell(1, 0).text = "❌ Toxic dispersants (Corexit)\n❌ High-pressure hot washing\n❌ 'Cleanup complete' fraud"
table.cell(1, 1).text = "✅ Bioremediation (5× degradation)\n✅ Nutrient optimization\n✅ Plant-based surfactants"
table.cell(1, 2).text = "Oil Pollution Act (1990):\n- Double-hull tankers\n- Spill liability funds\n- Response planning"

# Format table
for row in range(rows):
    for col in range(cols):
        cell = table.cell(row, col)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = (row == 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(240, 248, 255) if row % 2 == 0 else RGBColor(255, 255, 255)

# Speaker notes
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = ("SPEAKER NOTES:\n"
                                     "- Corexit toxicity: Killed organisms faster than oil\n"
                                     "- Bioremediation: Nitrogen optimization at 1-2 mg/L\n"
                                     "- Legal outcome: $5B → $500M settlement\n"
                                     "VISUAL: Comparison table with icons")

# Slide 7: Conclusion & Legacy
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Conclusion & Legacy"
content = slide.shapes.placeholders[1]
tf = content.text_frame
tf.text = "Unlearned Lessons:"
p = tf.add_paragraph()
p.text = "• 80% of Alaskan oil infrastructure remains aging"
p = tf.add_paragraph()
p.text = "• Arctic drilling expands despite risks"
p = tf.add_paragraph()
p.text = "\nPath Forward:"
p = tf.add_paragraph()
p.text = "• Prioritize non-toxic cleanup tech:"
p = tf.add_paragraph()
p.text = "  - Bio-aerogels, plant-based surfactants"
p = tf.add_paragraph()
p.text = "• Empower Indigenous-led monitoring"
p = tf.add_paragraph()
p.text = "• Enforce corporate accountability"
p = tf.add_paragraph()
p.text = "\n\"The spill was not an accident but a corporate crime\" - G. Palast (2008)"

# Format quote
for paragraph in tf.paragraphs:
    if "corporate crime" in paragraph.text:
        paragraph.font.italic = True
        paragraph.alignment = PP_ALIGN.CENTER

# Speaker notes
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = ("SPEAKER NOTES:\n"
                                     "- Highlight current risks: 23 aging tankers still in Alaskan waters\n"
                                     "- Mention promising tech: Nanocellulose aerogels absorb 50x weight in oil\n"
                                     "- Call to action: Support Indigenous conservation initiatives")

# Save presentation
prs.save('Exxon_Valdez_Oil_Spill_Presentation.pptx')
print("Presentation created successfully!")